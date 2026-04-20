import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

def set_font(p, size=18, bold=False, color=(255, 255, 255)):
    p.font.name = "Calibri"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)

def add_background(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(26, 60, 94)  # #1A3C5E
    bg.line.fill.background()
    # Move to back
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

def add_title(slide, text):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    set_font(p, size=28, bold=True)
    
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.33), Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(46, 134, 171)  # #2E86AB
    bar.line.fill.background()

def add_bullets(slide, bullets, top=1.8, left=0.5, width=12.33, font_size=20):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = "• " + b
        set_font(p, size=font_size)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

# SLIDE 1: Title Slide
s1 = prs.slides.add_slide(blank_layout)
add_background(s1)

tb = s1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(4))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "YOLOv8-Based Brain Tumour Detection System"
set_font(p, size=40, bold=True)
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "AI-Powered MRI Analysis for Early Diagnosis"
set_font(p2, size=24)
p2.alignment = PP_ALIGN.CENTER

tb2 = s1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(2.5))
tf2 = tb2.text_frame
p1 = tf2.paragraphs[0]
p1.text = "Team: Anuja | Gajalakshmi | Kamalasani"
set_font(p1, size=20)
p1.alignment = PP_ALIGN.CENTER

p2 = tf2.add_paragraph()
p2.text = "Guide: Mr. S. Mohanraj, AP/ECE"
set_font(p2, size=20)
p2.alignment = PP_ALIGN.CENTER

p3 = tf2.add_paragraph()
p3.text = "Institution: Manakula Vinayagar Institute of Technology (MVIT)"
set_font(p3, size=20)
p3.alignment = PP_ALIGN.CENTER

p4 = tf2.add_paragraph()
p4.text = "\nDeveloped by: TB Solutions — TORCH BEARER"
set_font(p4, size=18, bold=True, color=(46, 134, 171))
p4.alignment = PP_ALIGN.CENTER

# SLIDE 2: Abstract
s2 = prs.slides.add_slide(blank_layout)
add_background(s2)
add_title(s2, "Abstract")
add_bullets(s2, [
    "Problem: Manual MRI screening for brain tumours is time-consuming and heavily reliant on radiologist expertise.",
    "Solution: An automated, AI-driven diagnostic system for rapid and reliable detection from MRI scans.",
    "Model: YOLOv8s natively deployed for high-efficiency image classification and bounding box detection.",
    "Performance: The trained model achieved an exceptional mAP@0.5 of 0.993 across 4 classes.",
    "Deployment: Fully functional web application running publicly on Hugging Face Spaces for immediate clinical access."
], font_size=24)

# SLIDE 3: Problem Statement
s3 = prs.slides.add_slide(blank_layout)
add_background(s3)
add_title(s3, "Problem Statement")
add_bullets(s3, [
    "Manual MRI analysis is highly time-consuming and inherently error-prone due to fatigue.",
    "Radiologists face an overwhelmingly high workload, delaying critical preliminary screenings.",
    "Existing generic AI models are not domain-specific and produce inconsistent accuracy for clinical use.",
    "There is a severe lack of simple, accessible user interfaces bridging AI research and clinical practice.",
    "Critical Need: A fast, highly accurate, automated detection system accessible anywhere via web."
], font_size=24)

# SLIDE 4: Existing Systems vs Proposed System
s4 = prs.slides.add_slide(blank_layout)
add_background(s4)
add_title(s4, "Existing Systems vs Proposed System")

table_shape = s4.shapes.add_table(5, 2, Inches(1), Inches(2), Inches(11.33), Inches(4))
table = table_shape.table

cols = ["Existing Systems", "Proposed System"]
for col_idx, col_name in enumerate(cols):
    cell = table.cell(0, col_idx)
    cell.text = col_name
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(26, 60, 94)
    set_font(cell.text_frame.paragraphs[0], size=22, bold=True)

rows_data = [
    ("General-purpose basic CNNs/YOLO limits", "Custom fine-tuned YOLOv8s architecture"),
    ("Low inconsistent accuracy (70-85%)", "High accuracy (mAP@0.5 = 0.993)"),
    ("Command-line only (No clinical UI)", "User-friendly web interface via Gradio"),
    ("High compute constraints for execution", "Optimized cloud deployment on Hugging Face")
]

for r_idx, row in enumerate(rows_data):
    for c_idx, val in enumerate(row):
        cell = table.cell(r_idx + 1, c_idx)
        cell.text = val
        cell.fill.solid()
        if r_idx % 2 == 0:
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            set_font(cell.text_frame.paragraphs[0], size=18, color=(26, 60, 94))
        else:
            cell.fill.fore_color.rgb = RGBColor(235, 244, 250)
            set_font(cell.text_frame.paragraphs[0], size=18, color=(26, 60, 94))

# SLIDE 5: System Architecture
s5 = prs.slides.add_slide(blank_layout)
add_background(s5)
add_title(s5, "System Architecture")

steps = [
    "Stage 1: Input\n(MRI Image Upload)",
    "Stage 2: Preprocessing\n(Resize & Normalize)",
    "Stage 3: Inference\n(Detection & Classify)",
    "Stage 4: Output\n(Box, Label & Conf)"
]

left_pos = 0.5
for i, step in enumerate(steps):
    shape = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_pos), Inches(3), Inches(2.2), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(46, 134, 171)
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = step
    set_font(p, size=18, bold=True)
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    if i < 3:
        arr = s5.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(left_pos + 2.3), Inches(3.5), Inches(0.8), Inches(0.5))
        arr.fill.solid()
        arr.fill.fore_color.rgb = RGBColor(255, 255, 255)
        arr.line.fill.background()
    left_pos += 3.2

# SLIDE 6: Dataset & Training
s6 = prs.slides.add_slide(blank_layout)
add_background(s6)
add_title(s6, "Dataset & Training")
add_bullets(s6, [
    "Dataset: Brain Tumour MRI Dataset (Kaggle — masoudnickparvar)",
    "Total images: 8,269 splits (Train: 5,680 | Val: 1,717 | Test: 872)",
    "Classes: Glioma, Meningioma, Pituitary Adenoma, No Tumour",
    "Data Augmentation applied:",
    "   • Horizontal/Vertical Flip, Rotation ±15°",
    "   • HSV jitter, Mosaic 0.5 (Context retention)",
    "Training Environment: Google Colab (T4 GPU)",
    "Framework: Python, Ultralytics YOLOv8s, PyTorch 2.3.0",
    "Configuration: 23 epochs, Batch 16, Completion time ~2 hours"
], font_size=20)

# SLIDE 7: Training Results
s7 = prs.slides.add_slide(blank_layout)
add_background(s7)
add_title(s7, "Training Results")

table_shape = s7.shapes.add_table(7, 3, Inches(1), Inches(1.8), Inches(11.33), Inches(3.5))
table = table_shape.table

cols = ["Epoch", "mAP@0.5", "mAP@0.5:0.95"]
for col_idx, col_name in enumerate(cols):
    cell = table.cell(0, col_idx)
    cell.text = col_name
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(26, 60, 94)
    set_font(cell.text_frame.paragraphs[0], size=20, bold=True)

rows_data = [
    ("Epoch 1", "0.879", "0.816"), ("Epoch 5", "0.963", "0.911"),
    ("Epoch 10", "0.981", "0.961"), ("Epoch 15", "0.992", "0.877"),
    ("Epoch 20", "0.993", "0.988"), ("Epoch 23", "0.993", "0.981")
]

for r_idx, row in enumerate(rows_data):
    for c_idx, val in enumerate(row):
        cell = table.cell(r_idx + 1, c_idx)
        cell.text = val
        cell.fill.solid()
        if r_idx % 2 == 0:
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            set_font(cell.text_frame.paragraphs[0], size=18, color=(26, 60, 94))
        else:
            cell.fill.fore_color.rgb = RGBColor(235, 244, 250)
            set_font(cell.text_frame.paragraphs[0], size=18, color=(26, 60, 94))

box = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(5.8), Inches(11.33), Inches(1))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(235, 244, 250)
tf = box.text_frame
p = tf.paragraphs[0]
p.text = "FINAL METRICS: mAP@0.5 = 0.993  |  Precision = 0.977  |  Recall = 0.984"
set_font(p, size=24, bold=True, color=(26, 60, 94))
p.alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# SLIDE 8: Model Architecture
s8 = prs.slides.add_slide(blank_layout)
add_background(s8)
add_title(s8, "Model Architecture")
add_bullets(s8, [
    "YOLOv8s Specification Details:",
    "   • Total Layers: 225",
    "   • Parameters: 11.1 Million",
    "   • Processing Power: 28.7 GFLOPs",
    "   • Native Input Size: 640×640 pixels",
    "",
    "Framework Mechanics:",
    "   • Backbone: CSPDarknet paired with C2f feature blocks",
    "   • Head: Decoupled detection head for specialized bounding & class mapping",
    "   • Pretrained weights: Baseline initialized on COCO dataset",
    "   • Fine-tuned fully on Brain Tumour MRI specialized dataset"
], font_size=22)

# SLIDE 9: Web Application (Gradio UI)
s9 = prs.slides.add_slide(blank_layout)
add_background(s9)
add_title(s9, "Web Application (Gradio UI)")
add_bullets(s9, [
    "Application Interface Features:",
    "   • Drag-and-drop MRI image upload module (JPG / PNG)",
    "   • Real-time inference showing annotated bounding boxes",
    "   • Clear display of exact Class label & Confidence score overlay",
    "   • Interactive Detection Sensitivity tuning slider (0.10–0.90)",
    "   • Formatted textual detection summary panel",
    "   • Persistent embedded academic and medical use disclaimer",
    "",
    "Cloud Environment:",
    "   • Platform: Deployed publicly via Hugging Face Spaces",
    "   • Compute: Serverless CPU processing for remote global access"
], font_size=22)

# SLIDE 10: Results & Demo
s10 = prs.slides.add_slide(blank_layout)
add_background(s10)
add_title(s10, "Results & Trial Demonstrations")

cards = [
    "Card 1: Glioma detected\nConfidence: 94.2%",
    "Card 2: Meningioma detected\nConfidence: 87.6%",
    "Card 3: Pituitary detected\nConfidence: 91.8%",
    "Card 4: No Tumour\nConfidence: 96.3%"
]

txBox = s10.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Model reliably tested and proven on clinically unseen MRI scans from test split."
set_font(p, size=18, color=(200, 200, 200))

left_pos = 1
for i, card in enumerate(cards):
    shape = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_pos), Inches(2.5), Inches(2.5), Inches(3.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(235, 244, 250)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = card
    set_font(p, size=24, bold=True, color=(26, 60, 94))
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    left_pos += 3

# SLIDE 11: Conclusion & Future Work
s11 = prs.slides.add_slide(blank_layout)
add_background(s11)
add_title(s11, "Conclusion & Future Work")

box1 = s11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.8), Inches(12.33), Inches(2.2))
box1.fill.solid()
box1.fill.fore_color.rgb = RGBColor(46, 134, 171)
tf1 = box1.text_frame
pHeader1 = tf1.paragraphs[0]
pHeader1.text = "Conclusion\n"
set_font(pHeader1, size=20, bold=True)
for b in [
    "YOLOv8s achieved an outstanding mAP@0.5 of 0.993 on targeted MRI dataset.",
    "The bespoke Gradio web dashboard successfully deployed to Hugging Face Spaces.",
    "Substantially reduces manual radiologist workload providing rapid preliminary screening."
]:
    p = tf1.add_paragraph()
    p.text = "• " + b
    set_font(p, size=18)

box2 = s11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.5), Inches(12.33), Inches(2.2))
box2.fill.solid()
box2.fill.fore_color.rgb = RGBColor(46, 134, 171)
tf2 = box2.text_frame
pHeader2 = tf2.paragraphs[0]
pHeader2.text = "Future Work\n"
set_font(pHeader2, size=20, bold=True)
for b in [
    "Extend detection capabilities to complex 3D volumetric MRI scan segmentations.",
    "Direct architectural integration with localized hospital DICOM imaging storage systems.",
    "Develop a responsive native mobile application platform for remote tele-diagnosis."
]:
    p = tf2.add_paragraph()
    p.text = "• " + b
    set_font(p, size=18)

# SLIDE 12: References & Thank You
s12 = prs.slides.add_slide(blank_layout)
add_background(s12)
add_title(s12, "References & Acknowledgements")

refs = [
    "1. Ultralytics YOLOv8 Documentation — docs.ultralytics.com",
    "2. Brain Tumour MRI Dataset — Kaggle (masoudnickparvar)",
    "3. Research: doi: 10.3389/fonc.2025.1643208",
    "4. Research: doi: 10.1016/j.aej.2025.10.019",
    "5. Gradio UI Interface Documentation — gradio.app/docs"
]

txBox = s12.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(2))
tf = txBox.text_frame
for i, ref in enumerate(refs):
    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
    p.text = ref
    set_font(p, size=16)

box_ty = s12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.66), Inches(4.2), Inches(6), Inches(2.5))
box_ty.fill.solid()
box_ty.fill.fore_color.rgb = RGBColor(235, 244, 250)
tf_ty = box_ty.text_frame
p1 = tf_ty.paragraphs[0]
p1.text = "Thank You\n"
set_font(p1, size=40, bold=True, color=(26, 60, 94))
p1.alignment = PP_ALIGN.CENTER
p2 = tf_ty.add_paragraph()
p2.text = "Developed by TB Solutions — TORCH BEARER"
set_font(p2, size=20, bold=True, color=(46, 134, 171))
p2.alignment = PP_ALIGN.CENTER
p3 = tf_ty.add_paragraph()
p3.text = '"Lightning the future, Transforming the world."'
set_font(p3, size=16, color=(26, 60, 94))
p3.alignment = PP_ALIGN.CENTER
tf_ty.vertical_anchor = MSO_ANCHOR.MIDDLE

prs.save("brain_tumour_detection_presentation.pptx")
